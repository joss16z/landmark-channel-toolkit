from __future__ import annotations

import io
import re
from dataclasses import dataclass
from typing import Any, Optional

import pandas as pd
import streamlit as st
from pptx import Presentation

# Optional PDF support
try:
    import pdfplumber
except Exception:  # pragma: no cover
    pdfplumber = None

from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

st.set_page_config(page_title="Landmark Channel Toolkit", page_icon="🏢", layout="wide")

# =========================================================
# Normalisation helpers
# =========================================================

def clean_unit(value: Any) -> str:
    """Canonical unit key: G.02 -> G02, Unit4104 -> 4104, 1.06 -> 106."""
    if value is None or (isinstance(value, float) and pd.isna(value)):
        return ""
    s = str(value).strip().upper()
    s = re.sub(r"^UNIT\s*", "", s, flags=re.I)
    s = re.sub(r"\.0$", "", s)
    s = re.sub(r"[^A-Z0-9]", "", s)
    return s


def display_unit(value: Any) -> str:
    if value is None or (isinstance(value, float) and pd.isna(value)):
        return ""
    s = str(value).strip().upper()
    s = re.sub(r"^UNIT\s*", "", s, flags=re.I)
    s = re.sub(r"\.0$", "", s)
    return s


def money_number(value: Any) -> Optional[float]:
    if value is None or (isinstance(value, float) and pd.isna(value)):
        return None
    s = str(value).strip()
    if not s or not re.search(r"\d", s):
        return None
    raw = re.sub(r"[^0-9.\-]", "", s)
    if raw in {"", "-", "."}:
        return None
    try:
        return float(raw)
    except Exception:
        return None


def format_money(value: Any) -> str:
    n = money_number(value)
    if n is None:
        if value is None or (isinstance(value, float) and pd.isna(value)):
            return ""
        return str(value).replace("A$", "$").strip()
    return f"${n:,.0f}"


def same_money(a: Any, b: Any) -> bool:
    na = money_number(a)
    nb = money_number(b)
    if na is None and nb is None:
        return True
    if na is None or nb is None:
        return False
    return round(na) == round(nb)


def norm_int(value: Any) -> str:
    if value is None or (isinstance(value, float) and pd.isna(value)):
        return ""
    s = str(value).strip()
    if not s:
        return ""
    n = money_number(s)
    if n is not None:
        return str(int(round(n)))
    return re.sub(r"\.0$", "", s)


def area_text(internal: Any, external: Any) -> str:
    i = norm_int(internal)
    e = norm_int(external)
    if i and e:
        return f"{i} + {e}"
    return i or e or ""


def best_match(columns, candidates):
    lower = {str(c).lower().strip(): c for c in columns}
    for cand in candidates:
        key = cand.lower().strip()
        if key in lower:
            return lower[key]
    for cand in candidates:
        key = cand.lower().strip()
        for c in columns:
            if key in str(c).lower().strip():
                return c
    return None


def index_or_zero(columns, value):
    return columns.index(value) if value in columns else 0

# =========================================================
# Excel ingest
# =========================================================

@dataclass
class ExcelConfig:
    unit_col: str
    price_col: str
    status_col: Optional[str] = None
    internal_col: Optional[str] = None
    external_col: Optional[str] = None
    available_only: bool = True


def read_uploaded_bytes(uploaded_file) -> bytes:
    return uploaded_file.getvalue()


def get_sheet_names(uploaded_file) -> list[str]:
    data = read_uploaded_bytes(uploaded_file)
    try:
        xls = pd.ExcelFile(io.BytesIO(data))
        return xls.sheet_names
    except Exception:
        try:
            tables = pd.read_html(io.BytesIO(data))
            return [f"Table {i+1}" for i in range(len(tables))] or ["Sheet1"]
        except Exception:
            return ["Sheet1"]


def load_excel(uploaded_file, sheet_name: str) -> pd.DataFrame:
    data = read_uploaded_bytes(uploaded_file)
    try:
        df = pd.read_excel(io.BytesIO(data), sheet_name=sheet_name)
        return df.dropna(how="all")
    except Exception as e1:
        try:
            tables = pd.read_html(io.BytesIO(data))
            idx = int(sheet_name.split(" ")[1]) - 1 if sheet_name.startswith("Table ") else 0
            return tables[idx].dropna(how="all")
        except Exception:
            raise e1


def build_price_table(df: pd.DataFrame, cfg: ExcelConfig) -> pd.DataFrame:
    work = df.copy()
    work["__unit"] = work[cfg.unit_col].apply(clean_unit)
    work["Unit"] = work[cfg.unit_col].apply(display_unit)
    work["Excel Price"] = work[cfg.price_col].apply(format_money)
    work["__price_num"] = work[cfg.price_col].apply(money_number)

    if cfg.status_col:
        work["Status"] = work[cfg.status_col].fillna("").astype(str).str.strip()
    else:
        work["Status"] = ""

    if cfg.internal_col:
        work["Internal"] = work[cfg.internal_col].apply(norm_int)
    else:
        work["Internal"] = ""
    if cfg.external_col:
        work["External"] = work[cfg.external_col].apply(norm_int)
    else:
        work["External"] = ""
    work["Excel Area"] = [area_text(i, e) for i, e in zip(work["Internal"], work["External"])]

    if cfg.available_only and cfg.status_col:
        work = work[work["Status"].str.lower().eq("available")]

    work = work[work["__unit"].astype(bool)].drop_duplicates("__unit", keep="first")
    return work[["__unit", "Unit", "Excel Area", "Excel Price", "Status", "Internal", "External"]].copy()

# =========================================================
# Floorplate parsing: PPT + PDF
# =========================================================

# Unit formats: Unit G.02, Unit G02, Unit 201, Unit4104
UNIT_PATTERN = re.compile(r"\bUnit\s*([A-Z]{0,4}\.?\d{1,5}[A-Z]?)\b", flags=re.I)
PRICE_PATTERN = re.compile(r"\$\s?[0-9][0-9,]*(?:\.\d+)?")
AREA_PATTERN = re.compile(r"\b(\d{1,3})\s*\+\s*(\d{1,3})\b")


def extract_unit_blocks(text: str, source: str = "", page_or_slide: Any = "") -> list[dict[str, Any]]:
    """Extract Unit, nearest area and price from text after each Unit marker."""
    text = text or ""
    matches = list(UNIT_PATTERN.finditer(text))
    rows = []
    for i, m in enumerate(matches):
        raw_unit = m.group(1)
        unit_key = clean_unit(raw_unit)
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        segment = text[start:end]
        price_m = PRICE_PATTERN.search(segment)
        area_m = AREA_PATTERN.search(segment)
        rows.append({
            "__unit": unit_key,
            "Unit": display_unit(raw_unit),
            "Floorplate Area": f"{area_m.group(1)} + {area_m.group(2)}" if area_m else "",
            "Floorplate Price": format_money(price_m.group(0)) if price_m else "",
            "Source": source,
            "Page/Slide": page_or_slide,
            "Text Snippet": ("Unit " + raw_unit + segment[:120]).replace("\n", " | "),
        })
    return rows


def parse_ppt_floorplate(ppt_bytes: bytes) -> pd.DataFrame:
    prs = Presentation(io.BytesIO(ppt_bytes))
    rows = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text or ""
            for row in extract_unit_blocks(text, source=f"Slide {slide_idx} Shape {shape_idx}", page_or_slide=slide_idx):
                rows.append(row)
    df = pd.DataFrame(rows)
    if df.empty:
        return pd.DataFrame(columns=["__unit", "Unit", "Floorplate Area", "Floorplate Price", "Source", "Page/Slide", "Text Snippet"])
    return df.drop_duplicates("__unit", keep="first")


def parse_pdf_floorplate(pdf_bytes: bytes) -> pd.DataFrame:
    if pdfplumber is None:
        raise RuntimeError("pdfplumber is not installed. Please add pdfplumber to requirements.txt")
    rows = []
    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        for page_idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            for row in extract_unit_blocks(text, source=f"PDF Page {page_idx}", page_or_slide=page_idx):
                rows.append(row)
    df = pd.DataFrame(rows)
    if df.empty:
        return pd.DataFrame(columns=["__unit", "Unit", "Floorplate Area", "Floorplate Price", "Source", "Page/Slide", "Text Snippet"])
    return df.drop_duplicates("__unit", keep="first")


def parse_floorplate(uploaded_file) -> pd.DataFrame:
    name = uploaded_file.name.lower()
    data = uploaded_file.getvalue()
    if name.endswith(".pdf"):
        return parse_pdf_floorplate(data)
    if name.endswith(".pptx") or name.endswith(".ppt"):
        return parse_ppt_floorplate(data)
    raise ValueError("Unsupported floorplate file type")

# =========================================================
# Audit logic
# =========================================================

def sort_key_unit(u: str):
    s = clean_unit(u)
    prefix = re.sub(r"\d.*$", "", s)
    nums = re.findall(r"\d+", s)
    n = int(nums[0]) if nums else 999999
    return (prefix, n, s)


def generate_update_report(price_df: pd.DataFrame, fp_df: pd.DataFrame) -> pd.DataFrame:
    excel_map = {r["__unit"]: r for _, r in price_df.iterrows()}
    fp_map = {r["__unit"]: r for _, r in fp_df.iterrows()}
    excel_units = set(excel_map)
    fp_units = set(fp_map)
    rows = []

    for unit in sorted(excel_units & fp_units, key=sort_key_unit):
        e = excel_map[unit]
        f = fp_map[unit]
        price_changed = not same_money(e["Excel Price"], f["Floorplate Price"])
        area_changed = bool(e.get("Excel Area", "")) and bool(f.get("Floorplate Area", "")) and str(e.get("Excel Area", "")).strip() != str(f.get("Floorplate Area", "")).strip()
        if price_changed and area_changed:
            result = "Price & Area Changed"
            action = "Update price and area"
        elif price_changed:
            result = "Price Changed"
            action = "Update price"
        elif area_changed:
            result = "Area Changed"
            action = "Update area"
        else:
            result = "Matched"
            action = "No action"
        rows.append({
            "Unit": e["Unit"] or f["Unit"],
            "Result": result,
            "Action": action,
            "Floorplate Price": f.get("Floorplate Price", ""),
            "Excel Price": e.get("Excel Price", ""),
            "Floorplate Area": f.get("Floorplate Area", ""),
            "Excel Area": e.get("Excel Area", ""),
            "Status": e.get("Status", ""),
            "Location": f.get("Source", ""),
            "Text Snippet": f.get("Text Snippet", ""),
        })

    for unit in sorted(excel_units - fp_units, key=sort_key_unit):
        e = excel_map[unit]
        rows.append({
            "Unit": e["Unit"],
            "Result": "Missing in Floorplate",
            "Action": "Add label manually if needed",
            "Floorplate Price": "",
            "Excel Price": e.get("Excel Price", ""),
            "Floorplate Area": "",
            "Excel Area": e.get("Excel Area", ""),
            "Status": e.get("Status", ""),
            "Location": "",
            "Text Snippet": "",
        })

    for unit in sorted(fp_units - excel_units, key=sort_key_unit):
        f = fp_map[unit]
        rows.append({
            "Unit": f["Unit"],
            "Result": "Extra in Floorplate",
            "Action": "Review/remove if not available",
            "Floorplate Price": f.get("Floorplate Price", ""),
            "Excel Price": "",
            "Floorplate Area": f.get("Floorplate Area", ""),
            "Excel Area": "",
            "Status": "",
            "Location": f.get("Source", ""),
            "Text Snippet": f.get("Text Snippet", ""),
        })

    report = pd.DataFrame(rows)
    if report.empty:
        return report
    order = {
        "Price & Area Changed": 1,
        "Price Changed": 2,
        "Area Changed": 3,
        "Missing in Floorplate": 4,
        "Extra in Floorplate": 5,
        "Matched": 9,
    }
    report["__order"] = report["Result"].map(order).fillna(99)
    report = report.sort_values(["__order", "Unit"], key=lambda col: col.map(lambda x: str(x)), kind="stable")
    # more natural secondary sort
    report = report.drop(columns="__order")
    return report

# =========================================================
# PPT price update keeping formatting where possible
# =========================================================

def replace_price_in_run_text(run_text: str, price_map: dict[str, str]) -> tuple[str, bool]:
    original = run_text
    matches = list(UNIT_PATTERN.finditer(run_text or ""))
    if not matches:
        return run_text, False
    pieces = []
    last = 0
    changed = False
    for i, m in enumerate(matches):
        unit = clean_unit(m.group(1))
        seg_start = m.end()
        seg_end = matches[i + 1].start() if i + 1 < len(matches) else len(run_text)
        pieces.append(run_text[last:seg_start])
        segment = run_text[seg_start:seg_end]
        if unit in price_map:
            new_price = price_map[unit]
            def repl_price(pm):
                nonlocal changed
                if not same_money(pm.group(0), new_price):
                    changed = True
                return new_price
            segment = PRICE_PATTERN.sub(repl_price, segment, count=1)
        pieces.append(segment)
        last = seg_end
    pieces.append(run_text[last:])
    return "".join(pieces), changed or ("".join(pieces) != original)


def update_ppt_prices_keep_format(ppt_bytes: bytes, price_df: pd.DataFrame) -> tuple[bytes, pd.DataFrame]:
    prs = Presentation(io.BytesIO(ppt_bytes))
    price_map = {r["__unit"]: r["Excel Price"] for _, r in price_df.iterrows()}

    before_df = parse_ppt_floorplate(ppt_bytes)
    rows = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if not getattr(shape, "has_text_frame", False):
                continue
            before_text = shape.text_frame.text or ""
            if not UNIT_PATTERN.search(before_text):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    new_text, _ = replace_price_in_run_text(run.text, price_map)
                    if new_text != run.text:
                        run.text = new_text

    out = io.BytesIO()
    prs.save(out)
    after_df = parse_ppt_floorplate(out.getvalue())
    report = generate_update_report(price_df, after_df)
    report["Update Note"] = report["Result"].map(lambda x: "PPT generated; no missing labels added" if x != "Matched" else "Checked")
    return out.getvalue(), report

# =========================================================
# Excel report formatting
# =========================================================

COLORS = {
    "header": "1F4E78",
    "white": "FFFFFF",
    "price": "FCE4D6",
    "area": "E2F0D9",
    "missing": "FFF2CC",
    "extra": "F4CCCC",
    "matched": "D9EAD3",
    "summary": "DDEBF7",
}


def style_worksheet(ws):
    thin = Side(style="thin", color="D9EAD3")
    for cell in ws[1]:
        cell.fill = PatternFill("solid", fgColor=COLORS["header"])
        cell.font = Font(color=COLORS["white"], bold=True)
        cell.alignment = Alignment(horizontal="center", vertical="center")
    ws.freeze_panes = "A2"
    ws.auto_filter.ref = ws.dimensions
    for row in ws.iter_rows():
        for cell in row:
            cell.border = Border(bottom=Side(style="thin", color="E6E6E6"))
            cell.alignment = Alignment(vertical="top", wrap_text=True)
    for col_cells in ws.columns:
        max_len = 0
        col_letter = get_column_letter(col_cells[0].column)
        for cell in col_cells:
            if cell.value is not None:
                max_len = max(max_len, len(str(cell.value)))
        ws.column_dimensions[col_letter].width = min(max(max_len + 2, 12), 48)


def result_fill(result: str):
    if result in ["Price Changed", "Price & Area Changed"]:
        return PatternFill("solid", fgColor=COLORS["price"])
    if result == "Area Changed":
        return PatternFill("solid", fgColor=COLORS["area"])
    if result == "Missing in Floorplate":
        return PatternFill("solid", fgColor=COLORS["missing"])
    if result == "Extra in Floorplate":
        return PatternFill("solid", fgColor=COLORS["extra"])
    if result == "Matched":
        return PatternFill("solid", fgColor=COLORS["matched"])
    return None


def report_to_xlsx(report: pd.DataFrame, price_df: pd.DataFrame, fp_df: pd.DataFrame, project_name: str) -> bytes:
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        counts = report["Result"].value_counts().to_dict() if not report.empty else {}
        summary = pd.DataFrame([
            {"Metric": "Project", "Value": project_name},
            {"Metric": "Excel Available Units", "Value": len(price_df)},
            {"Metric": "Floorplate Labelled Units", "Value": len(fp_df)},
            {"Metric": "Price Changed", "Value": counts.get("Price Changed", 0) + counts.get("Price & Area Changed", 0)},
            {"Metric": "Area Changed", "Value": counts.get("Area Changed", 0) + counts.get("Price & Area Changed", 0)},
            {"Metric": "Missing in Floorplate", "Value": counts.get("Missing in Floorplate", 0)},
            {"Metric": "Extra in Floorplate", "Value": counts.get("Extra in Floorplate", 0)},
            {"Metric": "Matched", "Value": counts.get("Matched", 0)},
        ])
        summary.to_excel(writer, index=False, sheet_name="Summary")
        report.to_excel(writer, index=False, sheet_name="All Checks")

        sheets = [
            ("Price Changed", ["Price Changed", "Price & Area Changed"]),
            ("Missing in Floorplate", ["Missing in Floorplate"]),
            ("Extra in Floorplate", ["Extra in Floorplate"]),
            ("Area Changed", ["Area Changed", "Price & Area Changed"]),
            ("Matched", ["Matched"]),
        ]
        for name, results in sheets:
            subset = report[report["Result"].isin(results)] if not report.empty else pd.DataFrame()
            subset.to_excel(writer, index=False, sheet_name=name[:31])

        wb = writer.book
        for ws in wb.worksheets:
            style_worksheet(ws)
            # colour rows by Result
            headers = [cell.value for cell in ws[1]]
            if "Result" in headers:
                result_col = headers.index("Result") + 1
                for r in range(2, ws.max_row + 1):
                    fill = result_fill(str(ws.cell(r, result_col).value or ""))
                    if fill:
                        for c in range(1, ws.max_column + 1):
                            ws.cell(r, c).fill = fill
            if ws.title == "Summary":
                for r in range(2, ws.max_row + 1):
                    ws.cell(r, 1).font = Font(bold=True)
                    ws.cell(r, 1).fill = PatternFill("solid", fgColor=COLORS["summary"])
    return output.getvalue()

# =========================================================
# UI controls
# =========================================================

def excel_controls(key_prefix: str):
    excel_file = st.file_uploader("Upload Price List Excel", type=["xls", "xlsx"], key=f"{key_prefix}_excel")
    if not excel_file:
        return None, None, None

    sheet_names = get_sheet_names(excel_file)
    sheet_name = st.selectbox("Excel sheet", sheet_names, index=0, key=f"{key_prefix}_sheet")
    try:
        df = load_excel(excel_file, sheet_name)
    except Exception as e:
        st.error("Excel 读取失败。请尝试把文件另存为 .xlsx 后再上传，或确认 requirements.txt 已包含 xlrd。")
        st.exception(e)
        return None, None, None

    columns = list(df.columns)
    default_unit = best_match(columns, ["Unit Number", "Unit", "Apt #", "Apt", "Apartment", "Unit No"])
    default_price = best_match(columns, ["Contract Price", "Price", "List Price"])
    default_status = best_match(columns, ["Status", "Availability"])
    default_internal = best_match(columns, ["Internal", "Internal (sqm)", "Internal Area", "Internal Size"])
    default_external = best_match(columns, ["External", "External (sqm)", "External Area", "External Size", "Balcony"])

    c1, c2, c3 = st.columns(3)
    with c1:
        unit_col = st.selectbox("Unit column", columns, index=index_or_zero(columns, default_unit), key=f"{key_prefix}_unit_col")
    with c2:
        price_col = st.selectbox("Price column", columns, index=index_or_zero(columns, default_price), key=f"{key_prefix}_price_col")
    with c3:
        status_options = [None] + columns
        status_index = status_options.index(default_status) if default_status in status_options else 0
        status_col = st.selectbox("Status column", status_options, index=status_index, key=f"{key_prefix}_status_col")

    c4, c5, c6 = st.columns(3)
    with c4:
        internal_options = [None] + columns
        internal_index = internal_options.index(default_internal) if default_internal in internal_options else 0
        internal_col = st.selectbox("Internal area column", internal_options, index=internal_index, key=f"{key_prefix}_internal_col")
    with c5:
        external_options = [None] + columns
        external_index = external_options.index(default_external) if default_external in external_options else 0
        external_col = st.selectbox("External area column", external_options, index=external_index, key=f"{key_prefix}_external_col")
    with c6:
        available_only = st.checkbox("Only use Available units from Excel", value=True, key=f"{key_prefix}_available_only")

    cfg = ExcelConfig(unit_col=unit_col, price_col=price_col, status_col=status_col, internal_col=internal_col, external_col=external_col, available_only=available_only)
    price_df = build_price_table(df, cfg)

    with st.expander("Excel preview"):
        st.dataframe(df.head(20), use_container_width=True)
    return df, price_df, cfg


def show_metrics(report: pd.DataFrame, price_df: pd.DataFrame, fp_df: pd.DataFrame):
    counts = report["Result"].value_counts().to_dict() if not report.empty else {}
    m1, m2, m3, m4, m5 = st.columns(5)
    m1.metric("Excel Available", len(price_df))
    m2.metric("Floorplate Units", len(fp_df))
    m3.metric("Price Changed", counts.get("Price Changed", 0) + counts.get("Price & Area Changed", 0))
    m4.metric("Missing", counts.get("Missing in Floorplate", 0))
    m5.metric("Extra", counts.get("Extra in Floorplate", 0))

# =========================================================
# App
# =========================================================

st.sidebar.title("功能")
page = st.sidebar.radio(
    "选择功能",
    ["📋 Floorplate Update Report", "✅ PPT价格更新", "🏗️ Floorplate Generator", "📊 Price Compare", "📦 Agent Package Generator"],
)

st.title("🏢 Landmark Channel Toolkit")
st.caption("V4 — Better Floorplate Update Report + PPT/PDF support")

if page == "📋 Floorplate Update Report":
    st.header("📋 Floorplate Update Report")
    st.write("上传 Price List Excel 和 Floorplate PPT/PDF，生成更清晰的核对报告：Price Changed / Missing / Extra / Area Changed。不会修改原文件。")

    c1, c2 = st.columns(2)
    with c1:
        floorplate_file = st.file_uploader("Upload Floorplate PPT/PDF", type=["pptx", "ppt", "pdf"], key="report_floorplate")
    with c2:
        project_name = st.text_input("Project Name", value="Floorplate", key="report_project")

    df, price_df, cfg = excel_controls("report")

    if floorplate_file and price_df is not None and st.button("Generate Floorplate Update Report", type="primary"):
        try:
            fp_df = parse_floorplate(floorplate_file)
            report = generate_update_report(price_df, fp_df)
            report_xlsx = report_to_xlsx(report, price_df, fp_df, project_name=project_name)
            st.success("Floorplate Update Report generated.")
            show_metrics(report, price_df, fp_df)
            st.download_button(
                "Download Floorplate Update Report",
                report_xlsx,
                file_name=f"{project_name.replace(' ', '_')}_Floorplate_Update_Report.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
            st.dataframe(report, use_container_width=True)
        except Exception as e:
            st.error("生成报告失败。请检查上传文件是否为可读取的 PPT/PDF 和 Excel。")
            st.exception(e)

elif page == "✅ PPT价格更新":
    st.header("✅ 根据 Excel 自动更新 PPT 里的 Unit 价格")
    st.write("只替换已存在的 `Unit xxx - $xxx` 里的价格；不新增 Missing 户型框；同时生成 Floorplate Update Report。")
    st.warning("为尽量不改变字体和格式，本功能只做价格字符串替换。复杂分段文本可能无法替换，会体现在报告中。")

    c1, c2 = st.columns(2)
    with c1:
        ppt_file = st.file_uploader("Upload Floorplate PPT / PPTX", type=["pptx", "ppt"], key="update_ppt")
    with c2:
        project_name = st.text_input("Project Name", value="Floorplate", key="update_project")

    df, price_df, cfg = excel_controls("update")

    if ppt_file and price_df is not None and st.button("Update PPT Prices", type="primary"):
        try:
            updated_ppt, report = update_ppt_prices_keep_format(ppt_file.getvalue(), price_df)
            after_df = parse_ppt_floorplate(updated_ppt)
            report_xlsx = report_to_xlsx(report, price_df, after_df, project_name=project_name)
            st.success("PPT updated and report generated.")
            show_metrics(report, price_df, after_df)
            c1, c2 = st.columns(2)
            with c1:
                st.download_button(
                    "Download Updated PPT",
                    updated_ppt,
                    file_name=f"{project_name.replace(' ', '_')}_Updated_Prices.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            with c2:
                st.download_button(
                    "Download Floorplate Update Report",
                    report_xlsx,
                    file_name=f"{project_name.replace(' ', '_')}_Floorplate_Update_Report.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            st.dataframe(report, use_container_width=True)
        except Exception as e:
            st.error("更新 PPT 失败。")
            st.exception(e)

else:
    st.header(page)
    st.info("This module is planned. Current active modules: Floorplate Update Report and PPT价格更新.")
